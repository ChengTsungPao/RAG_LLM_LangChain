import os
import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.docstore.document import Document

def load_all_documents_from_folder(folder_path):
    documents = []
    for root, _, files in os.walk(folder_path):
        for filename in files:
            filepath = os.path.join(root, filename)
            ext = os.path.splitext(filename)[1].lower()

            try:
                if ext == ".pdf":
                    with fitz.open(filepath) as pdf:
                        text = "".join([page.get_text() for page in pdf])
                elif ext == ".docx":
                    doc = DocxDocument(filepath)
                    text = "\n".join([para.text for para in doc.paragraphs])
                elif ext == ".pptx":
                    prs = Presentation(filepath)
                    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                elif ext in [".txt", ".md", ".py"]:
                    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                        text = f.read()
                else:
                    continue

                documents.append(Document(page_content=text, metadata={"source": filepath}))
            except Exception as e:
                print(f"❌ 無法讀取：{filepath}，原因：{e}")
    return documents
